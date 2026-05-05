"""Apply in-place updates to aste_581_final_slides.pptx.

Non-destructive (preserve formatting / position / size):

  Slide 9:  swap the active-vs-fixed comparison image blob
  Slide 11: swap the q_acc sweep summary image blob
  Slide 13: update two text values in the fault scenarios table
            "Loose χ² gate" row → NIS 1.96 → 1.80, ‖r‖ 4.0e-5 → 3.7e-5
  Slide 13: insert FIG 13A (miss vs σ_pix) and FIG 13B (miss vs t_c) inside
            the panel placeholders.  Idempotent — re-running deletes any
            previously inserted FIG 13A/B picture and re-adds the latest.

Image swaps work by overwriting the bytes inside the embedded ImagePart —
position, size, and cropping on the slide stay exactly as authored.

Text updates work by mutating run.text on a single-run text shape, which
preserves the run's font, color, and bold/italic.

Usage:
    python scripts/_apply_slide_updates.py [--deck PATH] [--dry-run]
"""
from __future__ import annotations

import argparse
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from _common import repo_path


SLIDE_5_NEW_IMAGE = "reports/pipeline_diagram.png"
SLIDE_9_NEW_IMAGE = "results/active_tracking/07_active_tracking_fixed_vs_active_comparison.png"
SLIDE_11_NEW_IMAGE = "results/mc/fine_tune/06e_sweep_summary.png"

# Slide 11 — body paragraph + footer. Old text was written for the 5-point
# sweep (q ∈ [1e-11, 1e-9], 500 trials). New sweep is 11 points over
# q ∈ [1e-12, 1e-7] at n=1000.
SLIDE_11_TEXT_UPDATES = [
    ("Text 6",  "5.9 → 4.3 ",                    "5.8 → 3.0 "),
    ("Text 6",  "tightens; pass rate holds at ", "tightens; pass rate stays "),
    ("Text 6",  "99%",                           "≥ 90%"),
    ("Text 27", "500 TRIALS / POINT · CHOSEN Q_ACC = 1E-9",
                "1000 TRIALS / POINT · CHOSEN Q_ACC = 1E-9"),
]

SLIDE_13_TEXT_UPDATES = [
    ("Text 45", "1.96",   "1.80"),
    ("Text 46", "4.0e-5", "3.7e-5"),
    # FIG 13A panel footer chip — n_seeds bumped from 16 → 80
    ("Text 10", "TUNED · 16 SEEDS / POINT", "TUNED · 80 SEEDS / POINT"),
]

# Slide-13 panel geometries (inches), measured from the deck.
# Panel A (FIG 13A · MISS VS Σ_PIX) and Panel B (FIG 13B · MISS VS T_C)
# are 5.65 × 6.82 in rounded-rect placeholders starting at top=3.26 in.
# Useful image area sits below the FIG-N header strip (~3.55 in) and above
# the footer caption (~9.55 in).
SLIDE_13_FIGURE_INSERTS = [
    {
        "image": "results/mc/sensitivity_n80_proper/13a_miss_vs_sigma.png",
        "left": 1.05, "top": 3.62, "width": 5.55, "height": 5.95,
        "marker_name": "FIG13A_PICTURE",
    },
    {
        "image": "results/mc/sensitivity_n80_proper/13b_miss_vs_tc.png",
        "left": 6.95, "top": 3.62, "width": 5.55, "height": 5.95,
        "marker_name": "FIG13B_PICTURE",
    },
]


def _swap_image_blob(slide, new_path: Path) -> tuple[bool, str]:
    """Replace the bytes inside the slide's first PICTURE shape's image part."""
    new_blob = new_path.read_bytes()
    for sh in slide.shapes:
        if sh.shape_type != 13:  # PICTURE
            continue
        rid = sh._element.blip_rId
        image_part = sh.part.related_part(rid)
        old_size = len(image_part.blob)
        image_part._blob = new_blob
        return True, f"swapped {old_size} → {len(new_blob)} bytes"
    return False, "no PICTURE shape found"


def _insert_picture(slide, image_path: Path, *,
                    left: float, top: float, width: float, height: float,
                    marker_name: str) -> tuple[bool, str]:
    """Insert a Picture into the slide at the given position (inches).

    Idempotent: removes any prior shape that we previously inserted under
    the same marker_name, so re-runs cleanly replace stale figures.
    """
    spTree = slide.shapes._spTree
    removed = 0
    for sh in list(slide.shapes):
        if sh.name == marker_name:
            spTree.remove(sh._element)
            removed += 1

    pic = slide.shapes.add_picture(
        str(image_path),
        Inches(left), Inches(top),
        width=Inches(width), height=Inches(height),
    )
    pic.name = marker_name
    return True, (f"inserted at ({left:.2f}, {top:.2f}) "
                  f"size {width:.2f}×{height:.2f} in"
                  + (f"; replaced {removed} prior" if removed else ""))


def _update_text_run(slide, shape_name: str, expect: str, new: str) -> tuple[bool, str]:
    """Mutate the matching run of the named text shape, preserving its font.

    Matches via stripped comparison so callers can include whitespace in
    expect/new to preserve inter-word spacing across runs.  When a match
    is found, the run's leading and trailing whitespace pattern is kept
    even if the new value has a different one.

    Idempotent: if a run already matches the new value, returns success
    with no mutation.
    """
    expect_s = expect.strip()
    new_s    = new.strip()

    def _preserve_pad(orig: str, new_core: str) -> str:
        lead  = orig[:len(orig) - len(orig.lstrip())]
        trail = orig[len(orig.rstrip()):]
        return lead + new_core + trail

    for sh in slide.shapes:
        if sh.name != shape_name or not sh.has_text_frame:
            continue
        for para in sh.text_frame.paragraphs:
            # Walk each run individually, since formatting lives at the run
            # level — match only single runs that contain exactly expect_s.
            for run in para.runs:
                rt = run.text.strip()
                if rt == new_s:
                    return True, f"already {new_s!r} (no-op)"
                if rt == expect_s:
                    run.text = _preserve_pad(run.text, new_s)
                    return True, f"{expect_s!r} → {new_s!r}"
        return False, (f"shape found but no run matching {expect_s!r} or {new_s!r} "
                       f"(had: {sh.text_frame.text!r})")
    return False, f"no shape named {shape_name!r}"


def main() -> None:
    p = argparse.ArgumentParser()
    p.add_argument("--deck",
                   default="reports/aste_581_final_slides.pptx",
                   help="Path to the .pptx (relative to repo root).")
    p.add_argument("--dry-run", action="store_true",
                   help="Don't write anything; just report what would change.")
    args = p.parse_args()

    deck_path = repo_path(args.deck)
    print(f"\n▸ Deck: {deck_path}")
    if not deck_path.exists():
        raise SystemExit(f"Deck not found: {deck_path}")

    backup = deck_path.with_suffix(".pre_update.pptx")
    if not args.dry_run:
        shutil.copy2(deck_path, backup)
        print(f"  Backup → {backup}")

    pres = Presentation(deck_path)
    slides = list(pres.slides)

    # ── Slide 5: image swap (pipeline diagram) ───────────────────────────
    print("\n▸ Slide 5 (pipeline diagram)")
    img5 = repo_path(SLIDE_5_NEW_IMAGE)
    if not img5.exists():
        print(f"  ⚠️  source image missing: {img5}")
    else:
        ok, msg = _swap_image_blob(slides[4], img5)
        print(f"  {'✓' if ok else '✗'} {msg}")

    # ── Slide 9: image swap ──────────────────────────────────────────────
    print("\n▸ Slide 9 (active vs fixed)")
    img9 = repo_path(SLIDE_9_NEW_IMAGE)
    if not img9.exists():
        print(f"  ⚠️  source image missing: {img9}")
    else:
        ok, msg = _swap_image_blob(slides[8], img9)
        print(f"  {'✓' if ok else '✗'} {msg}")

    # ── Slide 11: image swap ─────────────────────────────────────────────
    print("\n▸ Slide 11 (q_acc sweep)")
    img11 = repo_path(SLIDE_11_NEW_IMAGE)
    if not img11.exists():
        print(f"  ⚠️  source image missing: {img11}")
    else:
        ok, msg = _swap_image_blob(slides[10], img11)
        print(f"  {'✓' if ok else '✗'} {msg}")

    # ── Slide 11: text updates (body + footer) ───────────────────────────
    print("\n▸ Slide 11 (body + footer text)")
    for shape_name, expect, new in SLIDE_11_TEXT_UPDATES:
        ok, msg = _update_text_run(slides[10], shape_name, expect, new)
        print(f"  {'✓' if ok else '✗'} {shape_name}: {msg}")

    # ── Slide 13: text updates ───────────────────────────────────────────
    print("\n▸ Slide 13 (fault scenarios — Loose χ² gate row)")
    for shape_name, expect, new in SLIDE_13_TEXT_UPDATES:
        ok, msg = _update_text_run(slides[12], shape_name, expect, new)
        print(f"  {'✓' if ok else '✗'} {shape_name}: {msg}")

    # ── Slide 13: figure inserts (FIG 13A, FIG 13B) ──────────────────────
    print("\n▸ Slide 13 (FIG 13A / FIG 13B insert)")
    for spec in SLIDE_13_FIGURE_INSERTS:
        img_path = repo_path(spec["image"])
        if not img_path.exists():
            print(f"  ⚠️  source image missing: {img_path}")
            continue
        ok, msg = _insert_picture(
            slides[12], img_path,
            left=spec["left"], top=spec["top"],
            width=spec["width"], height=spec["height"],
            marker_name=spec["marker_name"],
        )
        print(f"  {'✓' if ok else '✗'} {spec['marker_name']} ← {img_path.name}: {msg}")

    if args.dry_run:
        print("\n(dry-run; no file written)")
        return

    pres.save(deck_path)
    print(f"\n▸ Saved {deck_path}")


if __name__ == "__main__":
    main()
