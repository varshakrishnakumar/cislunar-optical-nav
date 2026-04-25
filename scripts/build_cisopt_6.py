"""Rebuild CisOpt-5.pptx → CisOpt-6.pptx with targeted edits.

Changes:
  S1  fix picture aspect-ratio distortion (slides 1, 5, 10, 11, 12, 13)
  S2  slide 8  — add 3σ position-envelope plot in bottom-right empty area
  S3  slide 11 — shrink miss hist, add ΔV inflation hist + trace(P) scatter
  S4  slide 12 — relabel footer "V&V · FILTER CONSISTENCY", add truth≡estimate line
  S5  slide 13 — replace three right-column text cards with 5-row robustness table
  S6  slide 15 — add feature-tracking metrics panel + concrete numbers
"""
from __future__ import annotations

import copy
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
SRC = ROOT / "reports" / "CisOpt-5.pptx"
DST = ROOT / "reports" / "CisOpt-6.pptx"

# Style constants (match existing deck)
FONT_MONO = "IBM Plex Mono"
FONT_BODY = "Geist"
C_BODY = RGBColor(0xC7, 0xC9, 0xCF)
C_LABEL = RGBColor(0x8A, 0x87, 0x75)
C_CYAN = RGBColor(0x5C, 0xE1, 0xE6)
C_AMBER = RGBColor(0xFF, 0xB5, 0x47)
C_RED = RGBColor(0xE5, 0x48, 0x4D)
C_MUTED = RGBColor(0x5F, 0x60, 0x65)
C_FRAME = RGBColor(0x2A, 0x2C, 0x31)

SLIDE_W_IN = 20.0


def native_ratio(path: Path) -> float:
    with Image.open(path) as im:
        return im.size[0] / im.size[1]


def fit_ratio_in_bbox(x: float, y: float, w: float, h: float, target_r: float) -> tuple[float, float, float, float]:
    """Return (x, y, w, h) sized so that w/h == target_r, centered in (x,y,w,h)."""
    cur_r = w / h
    if target_r > cur_r:
        # target wider than bbox → shrink height
        new_h = w / target_r
        new_y = y + (h - new_h) / 2
        return x, new_y, w, new_h
    else:
        new_w = h * target_r
        new_x = x + (w - new_w) / 2
        return new_x, y, new_w, h


def set_bbox_in(shape, xywh_in):
    x, y, w, h = xywh_in
    shape.left = Inches(x)
    shape.top = Inches(y)
    shape.width = Inches(w)
    shape.height = Inches(h)


def fix_picture_distortion(slide, picture_name: str, frame_name: str | None = None):
    """Fix picture aspect ratio and (optionally) resize a light-fill frame card to match."""
    for shp in slide.shapes:
        if shp.name == picture_name and shp.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                nw, nh = shp.image.size
            except Exception:
                return
            target_r = nw / nh
            x, y, w, h = (shp.left / 914400, shp.top / 914400, shp.width / 914400, shp.height / 914400)
            new = fit_ratio_in_bbox(x, y, w, h, target_r)
            set_bbox_in(shp, new)
            if frame_name:
                resize_frame_to_picture(slide, frame_name, new)
            return


def resize_frame_to_picture(slide, frame_name: str, pic_bbox_in, pad: float = 0.02):
    """Resize a background card to match a picture's bbox with a small padding."""
    x, y, w, h = pic_bbox_in
    for shp in slide.shapes:
        if shp.name == frame_name:
            set_bbox_in(shp, (x - pad, y - pad, w + 2 * pad, h + 2 * pad))
            return


def add_text(slide, xywh_in, text, *, font=FONT_MONO, size_pt=9, color=C_LABEL,
             align=PP_ALIGN.LEFT, bold=False):
    x, y, w, h = xywh_in
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size_pt)
    r.font.color.rgb = color
    r.font.bold = bold
    return tb


def add_caption(slide, xywh_in, text, color=C_LABEL):
    """Standard FIG caption label in mono."""
    return add_text(slide, xywh_in, text, font=FONT_MONO, size_pt=9, color=color)


def add_picture_fit(slide, image_path: Path, bbox_in):
    """Insert picture into bbox, preserving aspect ratio, centered within bbox."""
    x, y, w, h = bbox_in
    target_r = native_ratio(image_path)
    x2, y2, w2, h2 = fit_ratio_in_bbox(x, y, w, h, target_r)
    return slide.shapes.add_picture(str(image_path), Inches(x2), Inches(y2), width=Inches(w2), height=Inches(h2))


def remove_shape_by_name(slide, name: str):
    for shp in list(slide.shapes):
        if shp.name == name:
            sp = shp._element
            sp.getparent().remove(sp)
            return True
    return False


def remove_shapes_by_names(slide, names: list[str]):
    names = set(names)
    for shp in list(slide.shapes):
        if shp.name in names:
            sp = shp._element
            sp.getparent().remove(sp)


def replace_shape_text(slide, shape_name: str, new_text: str, *,
                       font=None, size_pt=None, color=None, bold=None):
    for shp in slide.shapes:
        if shp.name == shape_name and shp.has_text_frame:
            tf = shp.text_frame
            # Preserve first run's styling by default
            if tf.paragraphs and tf.paragraphs[0].runs:
                first_run = tf.paragraphs[0].runs[0]
                cur_font = first_run.font.name or FONT_MONO
                cur_size = first_run.font.size or Pt(9)
                try:
                    cur_color = first_run.font.color.rgb
                except Exception:
                    cur_color = C_LABEL
                cur_bold = first_run.font.bold
            else:
                cur_font, cur_size, cur_color, cur_bold = FONT_MONO, Pt(9), C_LABEL, None
            tf.clear()
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = new_text
            r.font.name = font or cur_font
            r.font.size = Pt(size_pt) if size_pt else cur_size
            try:
                r.font.color.rgb = color or cur_color
            except Exception:
                pass
            r.font.bold = cur_bold if bold is None else bold
            return True
    return False


# ------------------------------------------------------------------
# Main transforms
# ------------------------------------------------------------------
def fix_all_picture_distortions(prs):
    """Step S1: correct aspect ratios for pictures that are distorted >1.5%
    and resize their background 'photo cards' to match."""
    # (slide_idx_1based, picture_name, frame_name_or_None)
    targets = [
        (1, "Picture 55", None),
        (5, "Picture 29", None),
        (10, "Picture 23", "Shape 7"),
        # Slide 11 handled in edit_slide_11 (restructured layout)
        (12, "Picture 39", "Shape 9"),
        (13, "Picture 40", "Shape 9"),
        (13, "Picture 42", "Shape 16"),
    ]
    for idx, name, frame in targets:
        fix_picture_distortion(prs.slides[idx - 1], name, frame)


def edit_slide_8(prs):
    """S2 · Add 3σ envelope plot in bottom-right empty area of slide 8."""
    s = prs.slides[7]
    plot = ROOT / "results/diagnostics/06_ekf/estimate_tracking_baseline/plots/pos_3sigma.png"
    if not plot.exists():
        return

    # Empty area on right column: y ≈ 8.40 → 10.45, x = 10.25 → 19.25
    bbox = (10.25, 8.50, 9.00, 1.95)
    # Frame rectangle (subtle outline) — use textbox outline style of existing cards
    add_picture_fit(s, plot, bbox)
    # Small caption inside the new section
    add_text(s, (10.25, 8.25, 6.0, 0.22), "FIG 08 · 3σ POSITION ENVELOPE · 301 STEPS",
             font=FONT_MONO, size_pt=9, color=C_CYAN)
    add_text(s, (17.00, 8.25, 2.25, 0.22), "ESTIMATE-TRACKING", font=FONT_MONO, size_pt=9,
             color=C_LABEL, align=PP_ALIGN.RIGHT)


def edit_slide_11(prs):
    """S3 · Shrink miss hist, add ΔV hist + trace(P) scatter."""
    s = prs.slides[10]
    miss_img = ROOT / "results/mc/baseline_tuned/06c_hist_miss_ekf.png"
    dv_img = ROOT / "results/mc/baseline_tuned/06c_hist_dv_inflation_pct.png"
    scatter = ROOT / "results/mc/baseline_tuned/06c_scatter_traceP_vs_miss.png"

    # 1) Remove the placeholder text box '[ MC MISS HISTOGRAM ... ]' (Text 29)
    remove_shape_by_name(s, "Text 29")
    # The old single light-card frames don't fit the new 3-figure layout; drop them.
    remove_shapes_by_names(s, ["Shape 27", "Shape 28"])

    # 2) Resize Picture 51 (miss hist) to upper-left of right half, preserving ratio
    for shp in s.shapes:
        if shp.name == "Picture 51":
            x, y, w, h = 10.42, 1.64, 4.30, 0  # compute h from ratio
            try:
                nw, nh = shp.image.size
                h = w / (nw / nh)
            except Exception:
                h = 2.39
            set_bbox_in(shp, (x, y, w, h))
            break

    # 4) Move FIG 11 caption line (Text 31, Text 32) below the new figures but above stats
    #    Actually there's still room at y=8.27 — keep it, just edit the text
    replace_shape_text(s, "Text 31", "FIG 11 · MISS · ΔV INFLATION · TRACE(P)", color=C_CYAN)
    replace_shape_text(s, "Text 32", "three views of the same posterior", color=C_LABEL)

    # 5) Add ΔV inflation histogram next to miss hist
    if dv_img.exists():
        add_picture_fit(s, dv_img, (14.92, 1.64, 4.30, 2.40))
        add_text(s, (14.92, 4.06, 4.30, 0.22), "ΔV INFLATION %", font=FONT_MONO, size_pt=9, color=C_AMBER)
        add_text(s, (14.92, 4.26, 4.30, 0.22), "median +1.2% vs omniscient", font=FONT_BODY, size_pt=10, color=C_BODY)

    # Caption under miss hist
    add_text(s, (10.42, 4.06, 4.30, 0.22), "TERMINAL MISS NORM", font=FONT_MONO, size_pt=9, color=C_CYAN)
    add_text(s, (10.42, 4.26, 4.30, 0.22), "‖r_ekf − r_tgt‖ · CR3BP LU", font=FONT_BODY, size_pt=10, color=C_BODY)

    # 6) Add trace(P) vs miss scatter centered below, spanning right half
    if scatter.exists():
        # Available y: 4.70 → 8.20 (stats start at 8.81) = 3.50 tall
        # Ratio 1.636: if h=3.50 → w=5.72
        bbox = (10.42, 4.70, 8.71, 3.50)
        add_picture_fit(s, scatter, bbox)
        add_text(s, (10.42, 8.23, 8.71, 0.22), "FIG 11C · TRACE(P) → MISS · each dot = 1 trial",
                 font=FONT_MONO, size_pt=9, color=C_CYAN)


def edit_slide_12(prs):
    """S4 · Relabel footer as V&V, add truth≡estimate line."""
    s = prs.slides[11]
    # Footer label
    replace_shape_text(s, "Text 35", "V&V · FILTER CONSISTENCY")
    # Also the header section pill at top can keep "05 · RESULTS" but the main heading
    # is already "CONSISTENCY · NIS / NEES · QACC SWEEP" which is fine.
    # Add a truth≡estimate note in the SWEEP VERDICT card.
    # Shape 31 (card), Text 32 (header), Text 33 (body)
    replace_shape_text(
        s, "Text 33",
        "99% pass-rate is retained across qacc ∈ [1e-11, 1e-9]. NEES median shrinks "
        "from 5.9 → 4.3 as qacc tightens. V&V check: estimate-tracking ≡ truth-tracking "
        "to 6 decimals (6.135e-5 pos err) — active pointing introduces no bias."
    )


def edit_slide_13(prs):
    """S5 · Replace three right-column cards with 5-row robustness table."""
    s = prs.slides[12]

    # Drop the three right-side cards and their labels.
    # Σ SWEEP card:  Shape 22 + Text 23 + Text 24
    # Tc SWEEP card: Shape 25 + Text 26/27/28 (T C SWEEP) + Text 29
    # DETUNED card:  Shape 30 + Text 31/32/33 (DETUNED (Q ACC =1E-6)) + Text 34
    remove_shapes_by_names(
        s,
        [
            "Shape 22", "Text 23", "Text 24",
            "Shape 25", "Text 26", "Text 27", "Text 28", "Text 29",
            "Shape 30", "Text 31", "Text 32", "Text 33", "Text 34",
        ],
    )

    # Header
    add_text(s, (13.65, 2.87, 5.60, 0.24), "ROBUSTNESS · FAULT SCENARIOS",
             font=FONT_MONO, size_pt=10, color=C_CYAN, bold=True)
    add_text(s, (13.65, 3.12, 5.60, 0.22), "single-trial diagnostics · 301 steps",
             font=FONT_BODY, size_pt=9, color=C_LABEL)

    # Table
    rows = [
        ("SCENARIO", "VALID", "NIS", "FINAL |r|", "—"),
        ("Nominal", "1.00", "1.75", "6.1e-5", "✓"),
        ("4% dropout", "0.96", "1.79", "5.3e-5", "✓"),
        ("Pixel outliers", "1.00", "1.72", "7.3e-5", "✓"),
        ("Loose χ² gate", "1.00", "1.96", "4.0e-5", "✓"),
        ("1-step meas delay", "0.98", "1.72", "3.0e-3", "✗"),
    ]
    table_x, table_y = 13.65, 3.40
    col_w = [1.90, 0.70, 0.70, 1.10, 0.45]  # sum = 4.85
    row_h = 0.42
    n_rows = len(rows)

    # Background frame
    from pptx.enum.shapes import MSO_SHAPE
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(table_x - 0.05),
                            Inches(table_y - 0.05), Inches(sum(col_w) + 0.10),
                            Inches(row_h * n_rows + 0.10))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0x17, 0x17, 0x1A)
    bg.line.color.rgb = RGBColor(0x2A, 0x2C, 0x31)
    bg.line.width = Emu(6350)
    # Send background behind text
    bg._element.getparent().remove(bg._element)
    s.shapes._spTree.insert(2, bg._element)

    # Column separator line below header row
    sep = s.shapes.add_connector(1, Inches(table_x),
                                 Inches(table_y + row_h - 0.03),
                                 Inches(table_x + sum(col_w)),
                                 Inches(table_y + row_h - 0.03))
    sep.line.color.rgb = RGBColor(0x2A, 0x2C, 0x31)

    for r_idx, row in enumerate(rows):
        xx = table_x
        is_header = r_idx == 0
        is_fail = row[-1] == "✗"
        is_gray = row[0] == "Loose χ² gate"
        for c_idx, cell in enumerate(row):
            if is_header:
                color = C_LABEL
                font = FONT_MONO
                size = 8
                bold = True
            else:
                if c_idx == 0:
                    font = FONT_MONO
                    size = 10
                    color = C_BODY if not is_fail else RGBColor(0xE5, 0x48, 0x4D)
                    bold = False
                elif c_idx == 4:
                    font = FONT_MONO
                    size = 12
                    if cell == "✓":
                        color = C_CYAN
                    elif cell == "✗":
                        color = C_RED
                    else:
                        color = C_LABEL
                    bold = True
                else:
                    font = FONT_MONO
                    size = 10
                    color = C_BODY if not is_fail else C_RED
                    bold = False
            align = PP_ALIGN.LEFT if c_idx == 0 else PP_ALIGN.CENTER
            add_text(s, (xx + 0.08, table_y + r_idx * row_h + 0.10, col_w[c_idx] - 0.16, row_h - 0.14),
                     cell, font=font, size_pt=size, color=color, align=align, bold=bold)
            xx += col_w[c_idx]

    # Footnote below table
    add_text(
        s, (13.65, table_y + n_rows * row_h + 0.15, 5.60, 0.60),
        "Delay failure is the known limit. σ_pix and t_c sweeps (left) scan continuous stressors; "
        "this table scans discrete fault modes.",
        font=FONT_BODY, size_pt=9, color=C_LABEL,
    )


def edit_slide_15(prs):
    """S6 · Add feature tracking metrics panel + numbers."""
    s = prs.slides[14]
    metrics_img = ROOT / "results/demos/08_feature_tracking_metrics.png"
    if not metrics_img.exists():
        return

    # Existing layout: Left text column (0.75 → ~4), right video (8.88 → 19.25).
    # Insert a small metrics image below the left-column text bullets (y ≈ 8.6 → 10.2),
    # and add a concrete-numbers strip above it.
    add_text(s, (0.75, 8.50, 4.00, 0.22),
             "PROTOTYPE METRICS · 420 FRAMES", font=FONT_MONO, size_pt=9, color=C_CYAN, bold=True)

    # Three numbers
    stats = [
        ("DETECT", "100%", "420 / 420 frames"),
        ("CENTROID", "0.60 px", "median residual"),
        ("RADIUS", "1.1 px", "median |Δr|"),
    ]
    col_w = 4.00 / 3
    for i, (lbl, val, sub) in enumerate(stats):
        x = 0.75 + i * col_w
        add_text(s, (x, 8.80, col_w - 0.05, 0.20), lbl, font=FONT_MONO, size_pt=8, color=C_LABEL)
        add_text(s, (x, 9.02, col_w - 0.05, 0.34), val, font=FONT_MONO, size_pt=14, color=C_CYAN, bold=True)
        add_text(s, (x, 9.40, col_w - 0.05, 0.20), sub, font=FONT_BODY, size_pt=8, color=C_BODY)


def main():
    prs = Presentation(SRC)

    fix_all_picture_distortions(prs)
    edit_slide_8(prs)
    edit_slide_11(prs)
    edit_slide_12(prs)
    edit_slide_13(prs)
    edit_slide_15(prs)

    prs.save(DST)
    print(f"wrote {DST}")


if __name__ == "__main__":
    main()
